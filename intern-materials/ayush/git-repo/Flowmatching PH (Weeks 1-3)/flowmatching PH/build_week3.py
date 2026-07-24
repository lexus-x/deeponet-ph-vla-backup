#!/usr/bin/env python
"""
build_week3.py
==============
Generates all assets and assembles the 'Week 3 Progress Report' deck.
Re-runnable: reads result JSONs dynamically (so re-running after the P2 extra
perturbation evals finish auto-completes the P2 table/chart).
Output: Week3_progress_report.pptx (white background, vibrant colors).
"""
import json, os
from pathlib import Path
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

ROOT = Path("/home/user/Desktop/Ayush PH test")
AS = ROOT / "hugging face ckp/ppt_assets"; AS.mkdir(parents=True, exist_ok=True)
DATASETS = {"libero10": 101469, "spatial": 52970, "object": 66984}

# vibrant palette
BLUE, RED, GREEN, PURPLE, ORANGE, YELLOW = "#2E86DE", "#EE5253", "#10AC84", "#5F27CD", "#FF9F43", "#FECA57"

def load(p):
    p = Path(p); return json.load(open(p)) if p.exists() else None

# --------------------------------------------------------------- epochs/time
def epochs(s1, b1, s2, b2, frames):
    return (s1*b1 + s2*b2)/frames

RUNS = {
 "V1":            dict(ep=epochs(5000,8,30000,4,DATASETS["libero10"]), mins=82.6+79.7, note="LIBERO-10, batch 8/4"),
 "Strong-10":     dict(ep=epochs(5000,16,38000,16,DATASETS["libero10"]), mins=287.8+265.1, note="LIBERO-10, batch 16"),
 "Strong-Spatial":dict(ep=epochs(5000,16,20000,16,DATASETS["spatial"]), mins=144.4+148.4, note="LIBERO-Spatial, batch 16"),
 "P1-sweep":      dict(ep=epochs(5000,16,20000,16,DATASETS["spatial"]), mins=174.4+193.1+193.6+163.0+148.4, note="per λ ≈ 7.5 ep"),
 "P2-finetune":   dict(ep=epochs(0,16,10000,16,DATASETS["object"]), mins=87.2+84.9+95.8+77.0+98.1+73.7, note="continued finetune, per λ"),
}

# --------------------------------------------------------------- concept image
def concept_image():
    fig, ax = plt.subplots(figsize=(13, 6)); ax.axis("off"); ax.set_xlim(0,13); ax.set_ylim(0,6)
    ax.add_patch(FancyBboxPatch((0.3,0.3),12.4,5.4,boxstyle="round,pad=0.1",fc="#F7F9FC",ec=PURPLE,lw=3))
    ax.text(6.5,5.2,"Persistent Homology for Robust Vision-Language-Action Models",
            ha="center",fontsize=20,weight="bold",color=PURPLE)
    ax.text(6.5,4.5,"Matching the topology of predicted vs expert action trajectories",
            ha="center",fontsize=13,color="#555")
    # left: noisy trajectory point cloud
    rng=np.random.default_rng(1); pts=np.cumsum(rng.normal(0,1,(18,2)),0)
    pts=(pts-pts.min(0)); pts=pts/pts.max(0)*2.0
    ax.scatter(pts[:,0]+1.0,pts[:,1]+1.2,s=80,color=BLUE,zorder=3)
    ax.plot(pts[:,0]+1.0,pts[:,1]+1.2,color=BLUE,alpha=0.4)
    ax.text(2.0,0.8,"action chunk\n(point cloud)",ha="center",fontsize=11,weight="bold",color=BLUE)
    ax.add_patch(FancyArrowPatch((3.6,2.0),(5.2,2.0),arrowstyle="-|>",mutation_scale=22,lw=2.5,color=GREEN))
    ax.text(4.4,2.4,"PH",ha="center",fontsize=13,weight="bold",color=GREEN)
    # middle: persistence barcode
    for i,(b,d) in enumerate([(0.1,2.4),(0.3,1.8),(0.5,1.1),(0.8,0.9),(1.0,1.4)]):
        ax.plot([5.6+b,5.6+d],[1.2+i*0.4,1.2+i*0.4],lw=6,color=ORANGE,solid_capstyle="round")
    ax.text(6.6,0.8,"persistence\nbarcode",ha="center",fontsize=11,weight="bold",color=ORANGE)
    ax.add_patch(FancyArrowPatch((8.2,2.0),(9.6,2.0),arrowstyle="-|>",mutation_scale=22,lw=2.5,color=GREEN))
    # right: loss
    ax.add_patch(FancyBboxPatch((9.8,1.4),2.6,1.4,boxstyle="round,pad=0.1",fc=YELLOW,ec="black",lw=2))
    ax.text(11.1,2.1,r"$L_{total}=$"+"\n"+r"$L1 + \lambda\,L_{PH}$",ha="center",va="center",fontsize=13,weight="bold")
    fig.savefig(AS/"concept_title.png",bbox_inches="tight",dpi=150); plt.close(fig); print("concept_title.png")

# --------------------------------------------------------------- architecture
def architecture():
    fig, ax = plt.subplots(figsize=(13,7)); ax.axis("off"); ax.set_xlim(0,13); ax.set_ylim(0,7)
    def box(x,y,w,h,t,c): ax.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.05",fc=c,ec="black",lw=1.5)); ax.text(x+w/2,y+h/2,t,ha="center",va="center",fontsize=10,weight="bold")
    def ar(x1,y1,x2,y2): ax.add_patch(FancyArrowPatch((x1,y1),(x2,y2),arrowstyle="-|>",mutation_scale=18,lw=1.8,color="#333"))
    box(0.3,5.3,2.4,0.9,"Camera 1\n(agentview)","#CDE7FF"); box(0.3,4.1,2.4,0.9,"Camera 2\n(wrist)","#CDE7FF")
    box(0.3,2.9,2.4,0.9,"Robot state\n(8-dim)","#D7F7E6"); box(0.3,1.7,2.4,0.9,"Language\ninstruction","#FFE6C7")
    box(3.5,3.3,3.0,2.4,"SmolVLM-2\nVision-Language\nBackbone (~350M)","#E7D6FF")
    box(7.2,3.6,2.6,1.8,"Action Expert\n(flow matching,\n~100M)","#FFD3D3")
    box(10.3,3.8,2.4,1.4,"Predicted\naction chunk\n(50 x 7)","#FFF3B0")
    box(7.2,1.3,2.6,1.2,"L1 loss\n(predicted vs\nexpert actions)","#D5D5D5")
    box(10.3,1.3,2.4,1.2,"PH loss\n(topology of\naction chunk)","#FFB3B3")
    for y in (5.75,4.55,3.35,2.15): ar(2.7,y,3.5,4.5)
    ar(6.5,4.5,7.2,4.5); ar(9.8,4.5,10.3,4.5); ar(8.5,3.6,8.5,2.5); ar(11.5,3.8,11.5,2.5); ar(10.3,1.9,9.8,1.9)
    ax.text(6.5,0.6,r"$L_{total} = L1 + \lambda \cdot L_{PH}$    (λ = 0.1)",ha="center",fontsize=15,weight="bold",
            bbox=dict(boxstyle="round",fc=YELLOW,ec="black"))
    ax.text(6.5,6.6,"SmolVLA + Persistent-Homology (training-time only)",ha="center",fontsize=13,weight="bold",color=PURPLE)
    ax.text(12.5,6.6,"next →",ha="right",fontsize=12,style="italic",color=GREEN,weight="bold")
    fig.savefig(AS/"architecture.png",bbox_inches="tight",dpi=150); plt.close(fig); print("architecture.png")

# --------------------------------------------------------------- result bar chart
def result_bars(title, groups, path):
    # groups: list of (label, baseline, ph)
    fig, ax = plt.subplots(figsize=(8.5,4.8)); x=np.arange(len(groups)); w=0.36
    ax.bar(x-w/2,[g[1] for g in groups],w,label="flow-matching",color=BLUE)
    ax.bar(x+w/2,[g[2] for g in groups],w,label="flow + PH",color=RED)
    for i,g in enumerate(groups):
        ax.text(i-w/2,g[1]+1,f"{g[1]:.0f}",ha="center",fontsize=9,weight="bold")
        ax.text(i+w/2,g[2]+1,f"{g[2]:.0f}",ha="center",fontsize=9,weight="bold")
    ax.set_xticks(x); ax.set_xticklabels([g[0] for g in groups],fontsize=11)
    ax.set_ylim(0,100); ax.set_ylabel("success rate (%)"); ax.set_title(title,fontsize=13,weight="bold")
    ax.legend(); ax.grid(alpha=0.25,axis="y")
    fig.savefig(path,bbox_inches="tight",dpi=150); plt.close(fig)

# --------------------------------------------------------------- sweep table img
def table_image(title, headers, rows, best_idx, path):
    fig, ax = plt.subplots(figsize=(9.5, 0.6+0.5*len(rows))); ax.axis("off")
    tbl = ax.table(cellText=rows, colLabels=headers, loc="center", cellLoc="center")
    tbl.auto_set_font_size(False); tbl.set_fontsize(12); tbl.scale(1,1.6)
    for j in range(len(headers)):
        c=tbl[0,j]; c.set_facecolor(PURPLE); c.set_text_props(color="white",weight="bold")
    for i in range(len(rows)):
        for j in range(len(headers)):
            c=tbl[i+1,j]
            if i==best_idx: c.set_facecolor(YELLOW); c.set_text_props(weight="bold")
            elif i==0: c.set_facecolor("#EEF2FF")
            else: c.set_facecolor("#FFFFFF")
    ax.set_title(title,fontsize=14,weight="bold",color=PURPLE,pad=14)
    fig.savefig(path,bbox_inches="tight",dpi=150); plt.close(fig)

# --------------------------------------------------------------- sweep chart
def sweep_chart(title, lams, acc, rob, base_acc, base_rob, path, best_lam=None):
    fig,(a1,a2)=plt.subplots(1,2,figsize=(12,4.6))
    a1.plot(lams,acc,"o-",color=BLUE,lw=2,ms=8); a1.axhline(base_acc,ls="--",color="#888",label=f"baseline {base_acc:.0f}%")
    a1.set_xscale("log"); a1.set_title("λ  vs  Success rate (accuracy)",fontsize=12,weight="bold")
    a1.set_xlabel("λ (log)"); a1.set_ylabel("success rate (%)"); a1.set_ylim(0,100); a1.grid(alpha=0.3); a1.legend()
    a2.plot(lams,rob,"s-",color=RED,lw=2,ms=8); a2.axhline(base_rob,ls="--",color="#888",label=f"baseline {base_rob:.0f}%")
    a2.set_xscale("log"); a2.set_title("λ  vs  Robustness",fontsize=12,weight="bold")
    a2.set_xlabel("λ (log)"); a2.set_ylabel("robustness (%)"); a2.set_ylim(0,100); a2.grid(alpha=0.3); a2.legend()
    if best_lam:
        for ax_,arr in [(a1,acc),(a2,rob)]:
            bi=lams.index(best_lam); ax_.scatter([best_lam],[arr[bi]],s=260,facecolors="none",edgecolors=GREEN,lw=3,zorder=5)
    fig.suptitle(title,fontsize=14,weight="bold",color=PURPLE)
    fig.savefig(path,bbox_inches="tight",dpi=150); plt.close(fig)

# ====================================================== gather + generate
def pct(x): return round(x*100,1) if isinstance(x,(int,float)) else None

# P1 (spatial) data
def p1_data():
    strong=load("output/results_strong/success_rates.json")
    def get(node_ind,node_v):
        v=node_v; return (pct(node_ind["average"]), pct(v["viewpoint"]["average"]),
                          pct(v["lighting"]["average"]), pct(v["sensor_noise"]["average"]),
                          pct(v["average_over_perturbations"]))
    data={}
    data[0.0]=get(strong["LIBERO-SPATIAL"]["baseline"],strong["LIBERO-SPATIAL-V"]["baseline"])
    data[0.1]=get(strong["LIBERO-SPATIAL"]["ph"],strong["LIBERO-SPATIAL-V"]["ph"])
    for lam,t in [(0.02,"0p02"),(0.05,"0p05"),(0.2,"0p2"),(0.5,"0p5")]:
        d=load(f"output/results_sweep/lambda_{t}.json")
        if d: data[lam]=get(d["LIBERO-SPATIAL"]["ph"],d["LIBERO-SPATIAL-V"]["ph"])
    return dict(sorted(data.items()))

# P2 (object) data
def p2_data():
    HF="hugging face ckp/ph_object"; data={}
    for lam,t in [(0.0,"control"),(0.02,"lambda_0p02"),(0.05,"lambda_0p05"),(0.1,"lambda_0p1"),(0.2,"lambda_0p2"),(0.5,"lambda_0p5")]:
        d=load(f"{HF}/{t}/results/success_rates.json")
        if not d: continue
        ind=pct(d["LIBERO-OBJECT"]["ph"]["average"]); v=d["LIBERO-OBJECT-V"]["ph"]
        vp=pct(v.get("viewpoint",{}).get("average")); lg=pct(v.get("lighting",{}).get("average")); sn=pct(v.get("sensor_noise",{}).get("average"))
        avg=pct(v.get("average_over_perturbations"))
        data[lam]=(ind,vp,lg,sn,avg)
    return data

def best_lambda(data):  # max avg-robustness among λ>0 (fallback viewpoint)
    best=None;bv=-1
    for lam,(acc,vp,lg,sn,avg) in data.items():
        if lam==0: continue
        m=avg if avg is not None else vp
        if m is not None and m>bv: bv=m; best=lam
    return best

def fmt(x): return f"{x:.1f}" if isinstance(x,(int,float)) else "—"


# generate static assets
concept_image(); architecture()
# V1 + strong result bars
v1=load("output/v1/results/success_rates.json")
result_bars("V1 — LIBERO-10 (accuracy & robustness)",
    [("Accuracy",pct(v1["LIBERO-10"]["baseline"]["average"]),pct(v1["LIBERO-10"]["ph"]["average"])),
     ("Robustness",pct(v1["LIBERO-V"]["baseline"]["average_over_perturbations"]),pct(v1["LIBERO-V"]["ph"]["average_over_perturbations"]))],
    AS/"v1_results.png")
st=load("output/results_strong/success_rates.json")
result_bars("Stronger run — accuracy & robustness (λ = 0.1)",
    [("Obj-10 acc",pct(st["LIBERO-10"]["baseline"]["average"]),pct(st["LIBERO-10"]["ph"]["average"])),
     ("Obj-10 rob",pct(st["LIBERO-10-V"]["baseline"]["average_over_perturbations"]),pct(st["LIBERO-10-V"]["ph"]["average_over_perturbations"])),
     ("Spatial acc",pct(st["LIBERO-SPATIAL"]["baseline"]["average"]),pct(st["LIBERO-SPATIAL"]["ph"]["average"])),
     ("Spatial rob",pct(st["LIBERO-SPATIAL-V"]["baseline"]["average_over_perturbations"]),pct(st["LIBERO-SPATIAL-V"]["ph"]["average_over_perturbations"]))],
    AS/"strong_results.png")

# P1 sweep table + chart
p1=p1_data(); p1b=best_lambda(p1)
lams1=list(p1.keys()); rows1=[]
for lam in lams1:
    a,vp,lg,sn,avg=p1[lam]
    rows1.append([("baseline" if lam==0 else f"{lam}"),fmt(a),fmt(vp),fmt(lg),fmt(sn),fmt(avg)])
best1_idx=lams1.index(p1b)
table_image("LIBERO-Spatial — λ sweep (our trained model)",
            ["λ","Accuracy","Viewpoint-V","Lighting-V","Sensor-V","Robust avg"],rows1,best1_idx,AS/"p1_table.png")
sweep_chart("LIBERO-Spatial — λ sweep (our trained model)",
            [l for l in lams1 if l>0],[p1[l][0] for l in lams1 if l>0],[p1[l][4] for l in lams1 if l>0],
            p1[0.0][0],p1[0.0][4],AS/"p1_chart.png",best_lam=p1b)

# P2 sweep table + chart
p2=p2_data(); p2b=best_lambda(p2)
lams2=list(p2.keys()); rows2=[]
for lam in lams2:
    a,vp,lg,sn,avg=p2[lam]
    rows2.append([("baseline" if lam==0 else f"{lam}"),fmt(a),fmt(vp),fmt(lg),fmt(sn),fmt(avg)])
best2_idx=lams2.index(p2b) if p2b in lams2 else 0
table_image("LIBERO-Object — λ sweep (official checkpoint)",
            ["λ","Accuracy","Viewpoint-V","Lighting-V","Sensor-V","Robust avg"],rows2,best2_idx,AS/"p2_table.png")
robs2=[(p2[l][4] if p2[l][4] is not None else p2[l][1]) for l in lams2 if l>0]
sweep_chart("LIBERO-Object — λ sweep (official checkpoint)",
            [l for l in lams2 if l>0],[p2[l][0] for l in lams2 if l>0],robs2,
            p2[0.0][0],(p2[0.0][4] if p2[0.0][4] is not None else p2[0.0][1]),AS/"p2_chart.png",best_lam=p2b)
print("P1 best λ:",p1b," P2 best λ:",p2b)

# ====================================================== build PPT
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=13.333,7.5
NAVY=RGBColor(0x22,0x33,0x55); PUR=RGBColor(0x5F,0x27,0xCD); GRY=RGBColor(0x66,0x66,0x66)

def slide(): s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor(0xFF,0xFF,0xFF); return s
def title(s,t,size=30,color=PUR):
    tb=s.shapes.add_textbox(Inches(0.4),Inches(0.2),Inches(SW-0.8),Inches(1)); p=tb.text_frame.paragraphs[0]
    tb.text_frame.word_wrap=True; p.text=t; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(size); p.font.bold=True; p.font.color.rgb=color
def pic(s,img,top=1.15,maxh=5.5):
    if not Path(img).exists(): return
    iw,ih=Image.open(img).size; ar=iw/ih; h=maxh; w=h*ar
    if w>SW-1: w=SW-1; h=w/ar
    s.shapes.add_picture(str(img),Inches((SW-w)/2),Inches(top),width=Inches(w),height=Inches(h))
def caption(s,txt,color=GRY,size=14,bold=False):
    cb=s.shapes.add_textbox(Inches(0.5),Inches(SH-0.8),Inches(SW-1),Inches(0.6)); p=cb.text_frame.paragraphs[0]
    p.text=txt; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold
def two_pics(s,i1,i2,top=1.5,h=4.6):
    for k,img in enumerate([i1,i2]):
        if not Path(img).exists(): continue
        iw,ih=Image.open(img).size; ar=iw/ih; ww=h*ar
        if ww>SW/2-0.4: ww=SW/2-0.4; hh=ww/ar
        else: hh=h
        left=0.3+k*(SW/2-0.1)+((SW/2-0.4)-ww)/2
        s.shapes.add_picture(str(img),Inches(left),Inches(top),width=Inches(ww),height=Inches(hh))
def runtag(s,key):
    r=RUNS[key]; t=f"{r['ep']:.1f} epochs  •  {r['mins']:.0f} min total  ({r['mins']/60:.1f} h)"
    cb=s.shapes.add_textbox(Inches(0.5),Inches(SH-1.25),Inches(SW-1),Inches(0.5)); p=cb.text_frame.paragraphs[0]
    p.text=t; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(15); p.font.bold=True; p.font.color.rgb=RGBColor(0x10,0xAC,0x84)

# 0 TITLE
s=slide()
tb=s.shapes.add_textbox(Inches(1),Inches(2.1),Inches(SW-2),Inches(3)); tf=tb.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.text="Week 3 — Progress Report"; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(46); p.font.bold=True; p.font.color.rgb=PUR
p2=tf.add_paragraph(); p2.text="Persistent Homology regularization for SmolVLA on LIBERO"; p2.alignment=PP_ALIGN.CENTER; p2.font.size=Pt(20); p2.font.color.rgb=NAVY
p3=tf.add_paragraph(); p3.text="Ayush Shah"; p3.alignment=PP_ALIGN.CENTER; p3.font.size=Pt(26); p3.font.bold=True; p3.font.color.rgb=RGBColor(0xEE,0x52,0x53)

# 1 concept
s=slide(); title(s,"Concept — Persistent Homology (PH)"); pic(s,AS/"concept_title.png",top=1.3,maxh=5.0)
# 2 architecture
s=slide(); title(s,"Architecture — SmolVLA + PH"); pic(s,AS/"architecture.png",top=1.2,maxh=5.5)
caption(s,"Final loss:  L_total = L1 + λ · PH      (PH is training-time only)",color=NAVY,size=15,bold=True)
# 3 V1 what + results
s=slide(); title(s,"V1 — first run (LIBERO-10)")
pic(s,AS/"v1_results.png",top=1.3,maxh=4.6)
caption(s,"Parameters & inference latency stay identical (PH is train-only)",size=13)
runtag(s,"V1")
# 4 stronger run results
s=slide(); title(s,"Stronger run — results  (λ = 0.1)")
pic(s,AS/"strong_results.png",top=1.25,maxh=4.7)
caption(s,"Parameters & latency unchanged  •  LIBERO-10 + LIBERO-Spatial",size=13)
runtag(s,"Strong-10")
# 4b video
s=slide(); title(s,"Stronger run — PH succeeds where flow-matching fails")
if (AS/"ph_wins_task4.mp4").exists():
    iw,ih=Image.open(AS/"ph_wins_task4_frame.png").size if (AS/"ph_wins_task4_frame.png").exists() else (1032,290)
    ar=iw/ih; h=4.6; w=h*ar
    if w>SW-1: w=SW-1; h=w/ar
    s.shapes.add_movie(str(AS/"ph_wins_task4.mp4"),Inches((SW-w)/2),Inches(1.5),Inches(w),Inches(h),
                       poster_frame_image=str(AS/"ph_wins_task4_frame.png"))
caption(s,"LIBERO-Spatial task 4 — Left: flow-matching (FAIL)   |   Right: flow + PH (SUCCESS)",color=RGBColor(0xEE,0x52,0x53),size=15,bold=True)
# 5 motivation for sweep + P1 table
s=slide(); title(s,"Results not as required → sweep λ  (our trained model)")
pic(s,AS/"p1_table.png",top=1.4,maxh=4.6)
caption(s,f"Best λ highlighted = {p1b}   •   robustness on viewpoint / lighting / sensor-noise",size=13)
runtag(s,"P1-sweep")
# 6 P1 chart
s=slide(); title(s,"λ sweep — success vs robustness (our trained model)")
pic(s,AS/"p1_chart.png",top=1.4,maxh=4.8)
caption(s,f"green ring = best λ ({p1b})",size=13)
# 7 narrative for official ckpt
s=slide(); title(s,"Using an official online checkpoint")
tb=s.shapes.add_textbox(Inches(1.3),Inches(2.0),Inches(SW-2.6),Inches(3.5)); tf=tb.text_frame; tf.word_wrap=True
for i,b in enumerate(["We could not reach SmolVLA's maximum capacity on LIBERO with our own training.",
                      "So we used an official online checkpoint (lerobot/smolvla_libero) that claims higher accuracy.",
                      "We recreated the model from that checkpoint…",
                      "…and fine-tuned it with PH for different values of λ."]):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text="•  "+b; p.font.size=Pt(22); p.font.color.rgb=NAVY; p.space_after=Pt(12)
# 8 P2 table
s=slide(); title(s,"Official checkpoint — λ sweep table  (LIBERO-Object)")
pic(s,AS/"p2_table.png",top=1.4,maxh=4.6)
caption(s,f"Best λ highlighted = {p2b}",size=13)
runtag(s,"P2-finetune")
# 9 P2 chart
s=slide(); title(s,"Official checkpoint — success vs robustness")
pic(s,AS/"p2_chart.png",top=1.4,maxh=4.8)
caption(s,f"green ring = best λ ({p2b})",size=13)
# 10 thank you
s=slide()
tb=s.shapes.add_textbox(Inches(1),Inches(2.8),Inches(SW-2),Inches(2)); tf=tb.text_frame
p=tf.paragraphs[0]; p.text="Thank You"; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(54); p.font.bold=True; p.font.color.rgb=PUR
p2=tf.add_paragraph(); p2.text="Ayush Shah"; p2.alignment=PP_ALIGN.CENTER; p2.font.size=Pt(24); p2.font.color.rgb=RGBColor(0xEE,0x52,0x53)

OUT=ROOT/"Week3_progress_report.pptx"; prs.save(str(OUT))
print("saved",OUT)
