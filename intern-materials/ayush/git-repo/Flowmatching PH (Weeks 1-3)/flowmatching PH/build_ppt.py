#!/usr/bin/env python
"""
build_ppt.py
============
Assemble a simple, visual-heavy slide deck (college-student style) covering:
  * motivation + architecture diagram
  * PH method (loss equation, lambda)
  * v1 run  : success rate, robustness/perturbation, latency, parameters
  * stronger run : same set of plots (+ per-task win/loss)
  * a video where PH succeeds but flow-matching fails
Excludes the two in-progress (lambda-sweep / HF-PH) runs entirely.
Output: presentation.pptx in the project root.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

ROOT = Path("/home/user/Desktop/Ayush PH test")
V1 = ROOT / "output/v1/plots"
ST = ROOT / "output/plots_strong"
AS = ROOT / "hugging face ckp/ppt_assets"
OUTPPT = ROOT / "presentation.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xD6, 0x2C, 0x2C)
GREY = RGBColor(0x55, 0x55, 0x55)


def _title(slide, text, color=NAVY, size=30):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), SW - Inches(1), Inches(1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = color
    return tb


def _bg(slide, color=RGBColor(0xFF, 0xFF, 0xFF)):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


SW_IN, SH_IN = 13.333, 7.5


def img_slide(title, img, caption=None, top_in=1.2, max_h_in=5.6):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _title(s, title)
    if img and Path(img).exists():
        iw, ih = Image.open(img).size
        ar = iw / ih
        h_in = max_h_in; w_in = h_in * ar
        if w_in > SW_IN - 1:
            w_in = SW_IN - 1; h_in = w_in / ar
        left_in = (SW_IN - w_in) / 2
        s.shapes.add_picture(str(img), Inches(left_in), Inches(top_in),
                             width=Inches(w_in), height=Inches(h_in))
    if caption:
        cb = s.shapes.add_textbox(Inches(0.5), Inches(SH_IN - 0.85), Inches(SW_IN - 1), Inches(0.7))
        p = cb.text_frame.paragraphs[0]; p.text = caption; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14); p.font.color.rgb = GREY
    return s


def two_img_slide(title, img1, img2, caption=None):
    s = prs.slides.add_slide(BLANK); _bg(s); _title(s, title)
    half = SW_IN / 2
    for i, img in enumerate([img1, img2]):
        if img and Path(img).exists():
            iw, ih = Image.open(img).size; ar = iw / ih
            h_in = 4.8; w_in = h_in * ar
            if w_in > half - 0.4:
                w_in = half - 0.4; h_in = w_in / ar
            left_in = 0.3 + i * (half - 0.1) + (half - 0.4 - w_in) / 2
            s.shapes.add_picture(str(img), Inches(left_in), Inches(1.6),
                                 width=Inches(w_in), height=Inches(h_in))
    if caption:
        cb = s.shapes.add_textbox(Inches(0.5), SH - Inches(0.85), SW - Inches(1), Inches(0.7))
        p = cb.text_frame.paragraphs[0]; p.text = caption; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14); p.font.color.rgb = GREY
    return s


def bullet_slide(title, bullets, sub=None):
    s = prs.slides.add_slide(BLANK); _bg(s); _title(s, title)
    tb = s.shapes.add_textbox(Inches(1.2), Inches(1.8), SW - Inches(2.4), SH - Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + b; p.font.size = Pt(22); p.font.color.rgb = NAVY
        p.space_after = Pt(14)
    if sub:
        p = tf.add_paragraph(); p.text = sub; p.font.size = Pt(16); p.font.italic = True
        p.font.color.rgb = GREY
    return s


def section_slide(title, subtitle=None):
    s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
    tb = s.shapes.add_textbox(Inches(1), SH/2 - Inches(1), SW - Inches(2), Inches(2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle; p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xCF, 0xE8, 0xFF)
    return s


def video_slide(title, mp4, poster, caption=None):
    s = prs.slides.add_slide(BLANK); _bg(s); _title(s, title)
    if Path(mp4).exists():
        iw, ih = Image.open(poster).size if Path(poster).exists() else (1032, 290)
        ar = iw / ih
        h_in = 4.6; w_in = h_in * ar
        if w_in > SW_IN - 1: w_in = SW_IN - 1; h_in = w_in / ar
        left_in = (SW_IN - w_in) / 2
        s.shapes.add_movie(str(mp4), Inches(left_in), Inches(1.6), Inches(w_in), Inches(h_in),
                           poster_frame_image=str(poster) if Path(poster).exists() else None)
    if caption:
        cb = s.shapes.add_textbox(Inches(0.5), SH - Inches(0.9), SW - Inches(1), Inches(0.8))
        p = cb.text_frame.paragraphs[0]; p.text = caption; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = ACCENT
    return s


# ============================ SLIDES ============================
# 1 title
s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(1), Inches(2.3), SW - Inches(2), Inches(3)); tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Does Persistent Homology make SmolVLA more robust?"
p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
p2 = tf.add_paragraph(); p2.text = "Topological regularization for Vision-Language-Action models on LIBERO"
p2.alignment = PP_ALIGN.CENTER; p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xCF,0xE8,0xFF)
p3 = tf.add_paragraph(); p3.text = "Master's research project"
p3.alignment = PP_ALIGN.CENTER; p3.font.size = Pt(16); p3.font.color.rgb = RGBColor(0xCF,0xE8,0xFF)

# 2 motivation
bullet_slide("Motivation & Research Question", [
    "VLA models (SmolVLA) do well on clean benchmarks…",
    "…but can fail under visual perturbations (viewpoint, lighting, noise).",
    "Idea: add a Persistent-Homology (PH) loss that matches the topology of",
    "    predicted vs expert action trajectories.",
    "Question: does PH improve robustness without hurting accuracy or speed?",
])

# 3 architecture
img_slide("Architecture: SmolVLA + PH regularization", AS / "architecture.png",
          caption="PH is training-time only — inference is unchanged.")

# 4 method
bullet_slide("Method — PH loss", [
    "Total loss  =  flow-matching loss  +  λ · PH loss",
    "PH loss = difference in sorted top-k pairwise distances of action chunks",
    "Two-stage training: (1) warm-up head, (2) fine-tune full backbone",
    "bfloat16, AdamW, gradient checkpointing (Blackwell GPU)",
], sub="Regularization weight λ = 0.1")

# 5 section v1
section_slide("Experiment 1 — v1 run", "LIBERO-10 & LIBERO-V (visual perturbations)")
# 6-8 v1 plots
img_slide("v1 — Success rate (per task)", V1 / "libero10_per_task.png")
img_slide("v1 — Robustness to perturbations", V1 / "liberov_perturbation.png",
          caption="viewpoint / lighting / sensor-noise")
img_slide("v1 — Latency & Parameters", V1 / "params_latency.png",
          caption="PH adds zero inference cost (train-only) — latency & params identical")

# 9 section strong
section_slide("Experiment 2 — Stronger run", "Longer training (~6 epochs), 2 suites: LIBERO-10 + LIBERO-Spatial")
# 10 success
two_img_slide("Stronger run — Success rate per task",
              ST / "success_LIBERO-10_per_task.png", ST / "success_LIBERO-SPATIAL_per_task.png",
              caption="LIBERO-10 (left)  &  LIBERO-Spatial (right)")
# 11 robustness
two_img_slide("Stronger run — Robustness by perturbation",
              ST / "success_LIBERO-10-V_by_perturbation.png", ST / "success_LIBERO-SPATIAL-V_by_perturbation.png",
              caption="viewpoint / lighting / sensor-noise")
# 12 per-task win/loss
img_slide("Stronger run — Where PH helps (per-task Δ)", ST / "per_task_winloss_LIBERO-SPATIAL.png",
          caption="PH wins big on some tasks (e.g. task4), loses on others")
# 13 latency/params
img_slide("Stronger run — Latency & Parameters", ST / "params_latency.png",
          caption="Same architecture & latency for both — PH is free at inference")

# 14 video
section_slide("Qualitative result", "A task PH solves that flow-matching cannot")
video_slide("PH succeeds where flow-matching fails (LIBERO-Spatial, task 4)",
            AS / "ph_wins_task4.mp4", AS / "ph_wins_task4_frame.png",
            caption="Left: flow-matching (fails)   |   Right: flow + PH (succeeds)")

# 15 summary
bullet_slide("Summary", [
    "Built a full SmolVLA + PH training / eval / video pipeline on a Blackwell GPU",
    "PH is training-time only → zero inference cost (latency & params unchanged)",
    "Robustness vs accuracy shows a trade-off controlled by λ",
    "PH can solve specific tasks the baseline cannot (qualitative win)",
], sub="λ = 0.1 throughout")

prs.save(str(OUTPPT))
print(f"saved {OUTPPT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
